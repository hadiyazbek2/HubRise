"""
Generates HubRise Senior Project Presentation as a .pptx file (17 slides).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colours ───────────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # bright cyan
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCA, 0xF0, 0xF8)   # light cyan
GREY      = RGBColor(0xAA, 0xBB, 0xCC)
ORANGE    = RGBColor(0xFF, 0x8C, 0x00)   # warm accent

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank

def new_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, left, top, width, height,
          font_size=24, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tf_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tf_box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return tf_box

def accent_bar(slide, top=1.15, width=13.33, height=0.05):
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(top), Inches(width), Inches(height)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

def slide_number(slide, n):
    txbox(slide, str(n), 12.5, 7.1, 0.7, 0.35,
          font_size=11, color=GREY, align=PP_ALIGN.RIGHT)

def bullet_slide(slide, title, bullets, note=None):
    bg(slide)
    accent_bar(slide)
    txbox(slide, title, 0.4, 0.2, 12.5, 0.8,
          font_size=28, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.2), Inches(5.8))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"▸  {b}"
        run.font.size = Pt(20)
        run.font.name = "Arial"
        run.font.color.rgb = WHITE

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
bg(sl)
accent_bar(sl, top=3.55, width=13.33, height=0.06)

txbox(sl, "HubRise", 1, 1.2, 11, 1.4,
      font_size=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txbox(sl, "A Goal-Oriented Social Network for Android",
      1, 2.7, 11, 0.7, font_size=26, color=LIGHT, align=PP_ALIGN.CENTER)
txbox(sl, "Senior Project Presentation  |  2025 – 2026",
      1, 3.7, 11, 0.5, font_size=18, color=GREY, align=PP_ALIGN.CENTER)
txbox(sl, "Hadi Yazbek   |   Supervisor: Dr. Mohamad Rida Mortada",
      1, 4.3, 11, 0.5, font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "CS Engineering  |  github.com/hadiyazbek2/HubRise",
      1, 4.9, 11, 0.5, font_size=14, color=GREY, align=PP_ALIGN.CENTER)
slide_number(sl, 1)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Agenda
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
bullet_slide(sl, "Agenda", [
    "Problem Statement",
    "Project Idea & Concept",
    "Related Work & What Makes HubRise Different",
    "System Architecture",
    "Key Features Walkthrough",
    "Challenge System Deep Dive",
    "Code Highlights",
    "Technologies Used",
    "Limitations & Future Work",
    "Live Demo",
])
slide_number(sl, 2)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Problem Statement
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
bg(sl)
accent_bar(sl)
txbox(sl, "The Problem", 0.4, 0.2, 12.5, 0.8,
      font_size=28, bold=True, color=ACCENT)

txbox(sl, '"Social media gives you an audience, but not a purpose."',
      0.6, 1.3, 12, 0.8, font_size=22, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)

items = [
    ("Passive Consumption", "People scroll for entertainment — not growth."),
    ("No Accountability", "There is no social pressure to keep your commitments."),
    ("No Flexible Goals", "Fitness apps only track fitness. No app tracks ANY goal."),
    ("Shallow Community", "Following someone ≠ working toward something together."),
]
top = 2.3
for title, desc in items:
    txbox(sl, f"✕  {title}", 0.6, top, 5.0, 0.4, font_size=18, bold=True, color=ORANGE)
    txbox(sl, desc, 0.6, top + 0.38, 12, 0.4, font_size=16, color=WHITE)
    top += 0.9
slide_number(sl, 3)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — The Solution
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
bg(sl)
accent_bar(sl)
txbox(sl, "The Solution: HubRise", 0.4, 0.2, 12.5, 0.8,
      font_size=28, bold=True, color=ACCENT)

txbox(sl, "A social network where every interaction is tied to a goal.",
      0.6, 1.2, 12, 0.6, font_size=22, color=LIGHT, align=PP_ALIGN.CENTER)

pillars = [
    ("🏠  Hubs", "Join communities built around a shared goal — not just interests."),
    ("🎯  Challenges", "3 flexible models: stage, count, and streak. Cover any goal."),
    ("🤝  Accountability", "Peer validation + admin review = genuine, verified achievement."),
    ("📣  Social Layer", "Feed, likes, comments, follow/following — motivation through community."),
]
top = 2.1
for icon_title, desc in pillars:
    txbox(sl, icon_title, 0.6, top, 4.5, 0.45, font_size=20, bold=True, color=ACCENT)
    txbox(sl, desc, 0.6, top + 0.43, 12, 0.42, font_size=16, color=WHITE)
    top += 0.95
slide_number(sl, 4)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Related Work
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
bg(sl)
accent_bar(sl)
txbox(sl, "Related Work", 0.4, 0.2, 12.5, 0.8,
      font_size=28, bold=True, color=ACCENT)

headers = ["App", "Goal Types", "Community Feed", "Peer Validation", "Any Goal"]
col_w  = [2.2, 2.2, 2.6, 2.6, 2.2]
col_x  = [0.3, 2.55, 4.8, 7.45, 10.1]
rows = [
    ["Strava",    "Fitness only", "✓", "✗", "✗"],
    ["Habitica",  "Habits only",  "✗", "✗", "✗"],
    ["Instagram", "None",         "✓", "✗", "✗"],
    ["HubRise",   "Any",          "✓", "✓", "✓"],
]
top_start = 1.3
row_h = 0.7
for ci, header in enumerate(headers):
    txbox(sl, header, col_x[ci], top_start, col_w[ci], 0.5,
          font_size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    top_r = top_start + 0.55 + ri * row_h
    is_us = row[0] == "HubRise"
    rc = ACCENT if is_us else WHITE
    for ci, cell in enumerate(row):
        c = ACCENT if (is_us and ci > 0 and cell == "✓") else (
            RGBColor(0xFF, 0x4D, 0x4D) if cell == "✗" else rc
        )
        txbox(sl, cell, col_x[ci], top_r, col_w[ci], 0.55,
              font_size=16, bold=is_us, color=c, align=PP_ALIGN.CENTER)

txbox(sl, "HubRise is the only solution combining all four capabilities.",
      0.3, 6.3, 12.5, 0.5, font_size=16, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
slide_number(sl, 5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — System Architecture
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
bg(sl)
accent_bar(sl)
txbox(sl, "System Architecture", 0.4, 0.2, 12.5, 0.8,
      font_size=28, bold=True, color=ACCENT)

boxes = [
    (0.4, 1.4, 3.2, 1.6, "Android Client\n(Kotlin + MVVM)", ACCENT),
    (4.0, 1.4, 3.2, 1.6, "REST API\n(Django + DRF)", ORANGE),
    (7.6, 1.4, 3.2, 1.6, "Database\n(SQLite → PostgreSQL)", LIGHT),
    (1.2, 3.5, 2.4, 1.1, "Retrofit\n+ OkHttp", GREY),
    (4.0, 3.5, 2.4, 1.1, "JWT Auth\n(SimpleJWT)", GREY),
    (6.8, 3.5, 2.4, 1.1, "Media Storage\n(Filesystem)", GREY),
    (9.6, 3.5, 2.4, 1.1, "DataStore\n(Secure tokens)", GREY),
]
for bx, by, bw, bh, label, col in boxes:
    shape = sl.shapes.add_shape(
        1, Inches(bx), Inches(by), Inches(bw), Inches(bh)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x2D, 0x42)
    shape.line.color.rgb = col
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = col
    run.font.name = "Arial"

txbox(sl, "← JWT Bearer Token on every request →", 3.3, 2.05, 4.5, 0.4,
      font_size=12, color=GREY, align=PP_ALIGN.CENTER)
txbox(sl, "← SQL queries →", 7.1, 2.05, 3.2, 0.4,
      font_size=12, color=GREY, align=PP_ALIGN.CENTER)

txbox(sl,
      "25+ REST endpoints  •  MVVM on Android  •  JWT access (30 min) + refresh (7 days)",
      0.3, 5.1, 12.5, 0.5, font_size=15, color=LIGHT, align=PP_ALIGN.CENTER)
slide_number(sl, 6)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Authentication Feature
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
bullet_slide(sl, "Feature: User Registration & Authentication", [
    "3-step onboarding: email/password → username/name/DOB → interests/photo",
    "Real-time availability check for email and username during signup",
    "JWT access token (30 min TTL) + refresh token (7 days TTL)",
    "Tokens stored in Android DataStore — never SharedPreferences",
    "OkHttp Authenticator silently refreshes expired tokens (transparent to user)",
    "Refresh token blacklisted on logout — session cannot be reused",
    "Social login support: Google, Facebook, Apple",
])
slide_number(sl, 7)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Hubs & Feed Features
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
bullet_slide(sl, "Features: Feed & Hubs", [
    "Home Feed: personalized posts from joined Hubs, with likes & comments",
    "Global Feed: all public posts across the entire platform",
    "Hubs screen: 'My Hubs' and 'Explore Hubs' tabs (ViewPager2)",
    "Hub Detail: posts, challenges, member list, join/leave button",
    "Create Hub: name, description, cover image, category",
    "Hub Roles: Admin, Moderator, Member — each with different permissions",
    "Post types: Regular, Progress Update, Achievement Broadcast, Stage Proof, Count Entry, Streak Check-in",
])
slide_number(sl, 8)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Challenge System
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
bg(sl)
accent_bar(sl)
txbox(sl, "Feature: The Challenge System", 0.4, 0.2, 12.5, 0.8,
      font_size=28, bold=True, color=ACCENT)

models = [
    ("Stage-Based", "Sequential milestones.\nEach stage requires proof.\nVisual stage progress bar.", ACCENT),
    ("Count-Based", "Numeric target.\nLog entries (e.g. 100 km).\nAdditive or cumulative mode.", ORANGE),
    ("Streak-Based", "Daily/weekly check-ins.\nCalendar heatmap.\nGrace days for misses.", LIGHT),
]
for i, (title, desc, col) in enumerate(models):
    x = 0.4 + i * 4.3
    shape = sl.shapes.add_shape(1, Inches(x), Inches(1.3), Inches(4.0), Inches(3.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x2D, 0x42)
    shape.line.color.rgb = col
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = col
    r.font.name = "Arial"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "\n" + desc
    r2.font.size = Pt(14)
    r2.font.color.rgb = WHITE
    r2.font.name = "Arial"

completion_steps = "Completion Flow:  Member submits request  →  Peer micro-validation  →  Admin final review  →  Achievement post broadcast"
txbox(sl, completion_steps, 0.3, 4.55, 12.5, 0.6,
      font_size=15, color=LIGHT, align=PP_ALIGN.CENTER, italic=True)
txbox(sl, "Challenge Templates: 8 official templates across Fitness, Finance, Learning, Mindfulness, and more.",
      0.3, 5.3, 12.5, 0.5, font_size=14, color=GREY, align=PP_ALIGN.CENTER)
slide_number(sl, 9)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Social & Profile Features
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
bullet_slide(sl, "Features: Social, Profile & Notifications", [
    "Follow / Unfollow users — follower/following counts on profile",
    "Profile page: avatar, bio, post count, Hub count, wishlist URL",
    "Edit profile: name, bio, wishlist link",
    "Notifications: follows, likes, comments, challenge approvals/rejections",
    "Unread badge count on Notifications tab icon",
    "Comments: BottomSheet per post, live count shown on post card",
    "Explore screen: full-screen vertical scroll of public posts/videos",
    "Search: simultaneous search across Users, Hubs, Posts, Challenges",
])
slide_number(sl, 10)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Code Highlight 1: RetrofitClient
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
bg(sl)
accent_bar(sl)
txbox(sl, "Code Highlight #1 — RetrofitClient.kt", 0.4, 0.2, 12.5, 0.8,
      font_size=26, bold=True, color=ACCENT)
txbox(sl, "Transparent JWT Token Refresh", 0.4, 1.05, 12.5, 0.45,
      font_size=18, color=LIGHT, italic=True)

code = (
    "private val tokenAuthenticator = object : Authenticator {\n"
    "  override fun authenticate(route: Route?, response: Response): Request? {\n"
    "    // avoid infinite loop on refresh endpoint\n"
    "    if (response.request.url.encodedPath.contains(\"token/refresh\")) return null\n"
    "    val refreshToken = runBlocking { prefs.refreshToken.first() } ?: return navigateToLogin()\n"
    "    return try {\n"
    "      val refreshResponse = runBlocking { refreshService.refreshToken(...) }\n"
    "      runBlocking { prefs.saveAccessToken(refreshResponse.access) }\n"
    "      response.request.newBuilder().header(\"Authorization\", \"Bearer ${...}\").build()\n"
    "    } catch (e: Exception) { navigateToLogin(); null }\n"
    "  }\n"
    "}"
)
tb = slide.shapes.add_textbox(Inches(0.4), Inches(1.6), Inches(12.4), Inches(4.2)) if False else \
     sl.shapes.add_textbox(Inches(0.4), Inches(1.6), Inches(12.4), Inches(4.2))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = code
run.font.name = "Courier New"
run.font.size = Pt(12)
run.font.color.rgb = RGBColor(0xA8, 0xFF, 0x78)

txbox(sl,
      "When a request returns 401, OkHttp calls authenticate() automatically.\n"
      "A silent refresh happens — no error shown, original request retried with new token.",
      0.4, 5.95, 12.4, 0.8, font_size=15, color=WHITE)
slide_number(sl, 11)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Code Highlight 2: Challenge Model
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
bg(sl)
accent_bar(sl)
txbox(sl, "Code Highlight #2 — Challenge Model (Django)", 0.4, 0.2, 12.5, 0.8,
      font_size=26, bold=True, color=ACCENT)
txbox(sl, "Polymorphic Design — one model header, three configuration tables", 0.4, 1.05, 12.5, 0.45,
      font_size=18, color=LIGHT, italic=True)

code = (
    'class Challenge(models.Model):\n'
    '    MODEL_STAGE = "stage"  # → ChallengeStage rows\n'
    '    MODEL_COUNT = "count"  # → ChallengeCountConfig row\n'
    '    MODEL_STREAK = "streak" # → ChallengeStreakConfig row\n\n'
    '    hub = models.ForeignKey(Hub, ...)\n'
    '    progress_model = models.CharField(choices=MODEL_CHOICES)\n'
    '    is_main = models.BooleanField(default=False)\n'
    '    ends_at = models.DateTimeField(null=True, blank=True)\n\n'
    '# User progress per model:\n'
    '#   UserStageProgress  /  UserCountProgress  /  UserStreakProgress'
)
tb2 = sl.shapes.add_textbox(Inches(0.4), Inches(1.6), Inches(12.4), Inches(4.0))
tf2 = tb2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
run2 = p2.add_run()
run2.text = code
run2.font.name = "Courier New"
run2.font.size = Pt(13)
run2.font.color.rgb = RGBColor(0xA8, 0xFF, 0x78)

txbox(sl,
      "Single challenge row → model-specific config table.\n"
      "Clean schema, no nullable columns, supports thousands of concurrent participants.",
      0.4, 5.75, 12.4, 0.9, font_size=15, color=WHITE)
slide_number(sl, 12)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — Code Highlight 3: build_auth_response
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
bg(sl)
accent_bar(sl)
txbox(sl, "Code Highlight #3 — build_auth_response (Django)", 0.4, 0.2, 12.5, 0.8,
      font_size=26, bold=True, color=ACCENT)
txbox(sl, "Single helper — Signup, Login, Social Login all return the same response", 0.4, 1.05, 12.5, 0.45,
      font_size=18, color=LIGHT, italic=True)

code3 = (
    'def build_auth_response(user, request):\n'
    '    refresh = RefreshToken.for_user(user)\n'
    '    return {\n'
    '        "access_token":  str(refresh.access_token),\n'
    '        "refresh_token": str(refresh),\n'
    '        "user": UserSerializer(user, context={"request": request}).data,\n'
    '    }\n\n'
    '# Called identically from:\n'
    '#   SignupView.post()   →  return Response(build_auth_response(user, request), 201)\n'
    '#   LoginView.post()    →  return Response(build_auth_response(user, request), 200)\n'
    '#   SocialLoginView.post() → return Response(build_auth_response(user, request), 200)'
)
tb3 = sl.shapes.add_textbox(Inches(0.4), Inches(1.6), Inches(12.4), Inches(4.2))
tf3 = tb3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
run3 = p3.add_run()
run3.text = code3
run3.font.name = "Courier New"
run3.font.size = Pt(13)
run3.font.color.rgb = RGBColor(0xA8, 0xFF, 0x78)

txbox(sl,
      "DRY principle: one function, consistent response. "
      "Android client uses one AuthResponse data class for all three auth flows.",
      0.4, 5.95, 12.4, 0.8, font_size=15, color=WHITE)
slide_number(sl, 13)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — Technologies Used
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
bg(sl)
accent_bar(sl)
txbox(sl, "Technologies Used", 0.4, 0.2, 12.5, 0.8,
      font_size=28, bold=True, color=ACCENT)

backend = [
    "Python 3.12",
    "Django 6 + Django REST Framework",
    "SimpleJWT (JWT auth + token blacklisting)",
    "SQLite (dev) → PostgreSQL (prod)",
    "Pillow (image processing)",
]
android = [
    "Kotlin",
    "MVVM + Repository pattern",
    "Retrofit 2 + OkHttp 4",
    "Jetpack Navigation Component",
    "Android DataStore (Preferences)",
    "Glide (image loading)",
    "ViewPager2 + BottomSheet",
]
txbox(sl, "Backend", 0.4, 1.2, 5.5, 0.5, font_size=18, bold=True, color=ACCENT)
txbox(sl, "Android", 7.0, 1.2, 5.5, 0.5, font_size=18, bold=True, color=ORANGE)

top = 1.8
for item in backend:
    txbox(sl, f"▸  {item}", 0.5, top, 5.8, 0.45, font_size=16, color=WHITE)
    top += 0.46
top = 1.8
for item in android:
    txbox(sl, f"▸  {item}", 7.0, top, 5.8, 0.45, font_size=16, color=WHITE)
    top += 0.46
slide_number(sl, 14)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — Limitations
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
bullet_slide(sl, "Current Limitations", [
    "SQLite only — not suitable for concurrent production writes (PostgreSQL migration needed)",
    "No real-time push notifications (FCM) — notifications are pull-based",
    "Social login tokens not cryptographically verified — prototype only",
    "Media stored on local filesystem — needs cloud storage (S3) for production",
    "Single Django process — no Redis caching, no Celery task queue",
    "No comprehensive automated test suite",
])
slide_number(sl, 15)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16 — Future Work
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
bullet_slide(sl, "Future Work", [
    "Firebase Cloud Messaging (FCM) for real-time push notifications",
    "PostgreSQL migration + Redis caching layer",
    "AWS S3 / Google Cloud Storage for scalable media storage",
    "Achievement Badges displayed on user profiles",
    "Hub Leaderboards — rank members by challenge progress",
    "Direct Messaging between Hub members",
    "AI-powered challenge suggestions based on user interests",
    "iOS version (Swift / SwiftUI) consuming the same REST API",
])
slide_number(sl, 16)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 17 — Closing / Q&A
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
bg(sl)
accent_bar(sl, top=3.6, width=13.33, height=0.06)

txbox(sl, "Thank You", 1, 1.0, 11, 1.2,
      font_size=64, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txbox(sl, "Questions & Live Demo", 1, 2.4, 11, 0.7,
      font_size=28, color=LIGHT, align=PP_ALIGN.CENTER)
txbox(sl, "github.com/hadiyazbek2/HubRise", 1, 3.8, 11, 0.5,
      font_size=18, color=GREY, align=PP_ALIGN.CENTER)
txbox(sl, "Hadi Yazbek  —  CS Engineering  —  2025/2026", 1, 4.45, 11, 0.5,
      font_size=16, color=WHITE, align=PP_ALIGN.CENTER)
slide_number(sl, 17)

# ─────────────────────────────────────────────────────────────────────────────
output = "/home/hadi/AndroidStudioProjects/HubRise/HubRise_Presentation.pptx"
prs.save(output)
print(f"Presentation saved to: {output}")
